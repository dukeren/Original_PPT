import sys
import os
import openai
import logging
from PyQt5.QtWidgets import QApplication, QWidget, QRadioButton, QComboBox, QPushButton, QVBoxLayout, QHBoxLayout, QFileDialog, QMessageBox, QSpinBox, QLabel, QGroupBox, QStyleFactory
from PyQt5.QtCore import Qt
from PyQt5.QtGui import QPalette, QColor
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches
import requests
from PIL import Image
from io import BytesIO
import datetime
import random
import subprocess
import json
import time

# 配置日志
logging.basicConfig(level=logging.DEBUG, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

class PPTImageGenerator(QWidget):
    def __init__(self):
        super().__init__()
        self.config = self.load_config()
        QApplication.setStyle(QStyleFactory.create('Fusion'))
        self.initUI()
    
    def load_config(self):
        try:
            with open('config.json', 'r') as f:
                return json.load(f)
        except FileNotFoundError:
            logger.error("Configuration file not found. Using default settings.")
            return {}
        except json.JSONDecodeError:
            logger.error("Invalid JSON in configuration file. Using default settings.")
            return {}
        
    def initUI(self):
        mainLayout = QVBoxLayout()
        
        # 文件选择板块
        fileGroup = QGroupBox("文件选择")
        fileLayout = QHBoxLayout()
        self.filePathEdit = QLabel('未选择文件')
        self.fileSelectButton = QPushButton('选择文件')
        self.fileSelectButton.clicked.connect(self.selectFile)
        fileLayout.addWidget(self.filePathEdit)
        fileLayout.addWidget(self.fileSelectButton)
        fileGroup.setLayout(fileLayout)
        mainLayout.addWidget(fileGroup)

        # 生成范围板块
        rangeGroup = QGroupBox("生成范围")
        rangeLayout = QVBoxLayout()
        self.singleSlideRadio = QRadioButton('单页PPT图片生成')
        self.allSlidesRadio = QRadioButton('全部PPT图片生成')
        self.singleSlideRadio.setChecked(True)
        rangeLayout.addWidget(self.singleSlideRadio)
        rangeLayout.addWidget(self.allSlidesRadio)
        
        self.pageSelectLayout = QHBoxLayout()
        self.pageSelectLabel = QLabel('选择页面:')
        self.pageSelectSpinBox = QSpinBox()
        self.pageSelectSpinBox.setMinimum(1)
        self.pageSelectLayout.addWidget(self.pageSelectLabel)
        self.pageSelectLayout.addWidget(self.pageSelectSpinBox)
        rangeLayout.addLayout(self.pageSelectLayout)
        
        rangeGroup.setLayout(rangeLayout)
        mainLayout.addWidget(rangeGroup)

        # NLP Method selection
        nlpGroup = QGroupBox("自然语言识别模型")
        nlpLayout = QVBoxLayout()
        self.nlpMethod = QComboBox()
        self.nlpMethod.addItems(['Ollama', 'ChatGPT'])
        nlpLayout.addWidget(self.nlpMethod)
        nlpGroup.setLayout(nlpLayout)
        mainLayout.addWidget(nlpGroup)

        # 生成方式板块
        methodGroup = QGroupBox("图像生成方式")
        methodLayout = QVBoxLayout()
        self.generationMethod = QComboBox()
        self.generationMethod.addItems(['Comfyui生成', 'DALL-E', 'Pixabay', 'Unsplash'])
        methodLayout.addWidget(self.generationMethod)
        methodGroup.setLayout(methodLayout)
        mainLayout.addWidget(methodGroup)
        
        # 提交按钮
        self.submitButton = QPushButton('一键生成')
        self.submitButton.clicked.connect(self.onSubmit)
        mainLayout.addWidget(self.submitButton)
        
        self.setLayout(mainLayout)
        self.setWindowTitle('PPT图片生成 By 渡客')
        self.setGeometry(300, 300, 670, 300)

        # 连接单选按钮到更新UI的函数
        self.singleSlideRadio.toggled.connect(self.updatePageSelectVisibility)
        self.allSlidesRadio.toggled.connect(self.updatePageSelectVisibility)

        # 初始更新UI
        self.updatePageSelectVisibility()

    def updatePageSelectVisibility(self):
        is_single_slide = self.singleSlideRadio.isChecked()
        self.pageSelectLabel.setVisible(is_single_slide)
        self.pageSelectSpinBox.setVisible(is_single_slide)

    def selectFile(self):
        options = QFileDialog.Options()
        fileName, _ = QFileDialog.getOpenFileName(self, "选择PPT文件", "", "PowerPoint Files (*.pptx)", options=options)
        if fileName:
            self.filePathEdit.setText(fileName)
            # 更新页面选择的最大值
            prs = Presentation(fileName)
            self.pageSelectSpinBox.setMaximum(len(prs.slides))

    def onSubmit(self):
        logger.info("Submit button clicked")
        fileName = self.filePathEdit.text()
        if fileName and fileName != '未选择文件':
            logger.info(f"Selected file: {fileName}")
            isSingleSlide = self.singleSlideRadio.isChecked()
            generationMethod = self.generationMethod.currentText()
            pageNumber = self.pageSelectSpinBox.value() if isSingleSlide else None
            logger.info(f"Processing mode: {'Single Slide' if isSingleSlide else 'All Slides'}")
            logger.info(f"Generation method: {generationMethod}")
            if isSingleSlide:
                logger.info(f"Selected page number: {pageNumber}")
            self.processPPT(fileName, isSingleSlide, generationMethod, pageNumber)
        else:
            logger.warning("No file selected")
            QMessageBox.warning(self, "错误", "请先选择PPT文件")

    def processPPT(self, pptFile, isSingleSlide, generationMethod, pageNumber=None):
        logger.info(f"Processing PPT file: {pptFile}")
        prs = Presentation(pptFile)
        max_attempts = 3
        for attempt in range(max_attempts):
            if isSingleSlide:
                logger.info(f"Processing single slide: {pageNumber}")
                if 1 <= pageNumber <= len(prs.slides):
                    slide = prs.slides[pageNumber - 1]  # 调整为0基索引
                    self.processSlide(slide, generationMethod)
                else:
                    logger.warning(f"Invalid page number: {pageNumber}")
                    QMessageBox.warning(self, "错误", f"无效的页面号: {pageNumber}")
                    return
            else:
                logger.info("Processing all slides")
                for i, slide in enumerate(prs.slides):
                    logger.info(f"Processing slide {i+1}")
                    self.processSlide(slide, generationMethod)
            
            # 检查是否所有需要的更改都已完成
            if self.all_changes_completed():
                break
        else:
            logger.warning(f"Reached maximum attempts ({max_attempts}) for processing")
        
        # 创建输出目录
        output_dir = os.path.join(os.getcwd(), 'Outfile')
        os.makedirs(output_dir, exist_ok=True)
        
        # 生成带有随机时间后缀的文件名
        timestamp = datetime.datetime.now().strftime("%Y%m%d%H%M%S")
        random_suffix = str(random.randint(1000, 9999))
        output_filename = f'updated_{os.path.basename(pptFile)[:-5]}_{timestamp}_{random_suffix}.pptx'
        output_path = os.path.join(output_dir, output_filename)
        
        # 保存修改后的PPT
        logger.info(f"Saving updated PPT as: {output_path}")
        prs.save(output_path)
        
        # 打开生成文件目录
        self.open_output_directory(output_dir)
        
        QMessageBox.information(self, "完成", f"PPT处理完成!\n文件保存在: {output_path}")

    def open_output_directory(self, directory):
        if sys.platform == 'win32':
            os.startfile(directory)
        elif sys.platform == 'darwin':  # macOS
            subprocess.call(['open', directory])
        else:  # linux
            subprocess.call(['xdg-open', directory])

    def processSlide(self, slide, generationMethod):
        logger.info(f"Processing slide with layout: {slide.slide_layout.name}")
        shapes_to_process = list(slide.shapes)  # 创建一个副本
        for shape in shapes_to_process:
            logger.debug(f"Examining shape: {shape.name}, Type: {shape.shape_type}")
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or shape.name.startswith(('image', 'subimage')):
                logger.info(f"Found image shape: {shape.name}")
                content = self.findContentForImage(slide, shape)
                if content:
                    logger.info(f"Content found for image: {content[:250]}...")  # 只显示前500个字符
                    keywords = self.extractKeywords(content)
                    logger.info(f"Extracted keywords: {keywords}")
                    image = self.generateImage(keywords, generationMethod)
                    if image:
                        logger.info("Image generated successfully")
                        self.insertImage(slide, shape, image)
                    else:
                        logger.warning("Failed to generate image")
                else:
                    logger.warning(f"No content found for image: {shape.name}")

    def findContentForImage(self, slide, imageShape):
        image_name = imageShape.name
        content = ""
        logger.info(f"Finding content for image: {image_name}")

        if image_name.startswith('image'):
            number = image_name[5:]
            content_name = f'content{number}'
            logger.debug(f"Looking for content shape: {content_name}")
            content = self.findShapeTextByName(slide, content_name)
        
        elif image_name.startswith('subimage'):
            number = image_name[8:]
            subtitle_name = f'subtitle{number}'
            subcontent_name = f'subcontent{number}'
            logger.debug(f"Looking for subtitle shape: {subtitle_name}")
            logger.debug(f"Looking for subcontent shape: {subcontent_name}")
            subtitle = self.findShapeTextByName(slide, subtitle_name)
            subcontent = self.findShapeTextByName(slide, subcontent_name)
            content = f"{subtitle} {subcontent}".strip()

        # 如果没有找到对应的内容，尝试从其他占位符获取
        if not content:
            logger.warning("No specific content found, trying to find generic content")
            content = self.findGenericContent(slide)

        if content:
            logger.info(f"Content found: {content[:50]}...")  # 只显示前50个字符
        else:
            logger.warning("No content found, using default")
            content = "Generic image"

        return content

    def findShapeTextByName(self, slide, name):
        for shape in slide.shapes:
            if shape.name == name and hasattr(shape, 'text'):
                logger.debug(f"Found shape with name: {name}")
                return shape.text
        logger.debug(f"Shape not found: {name}")
        return ""

    def findGenericContent(self, slide):
        content = ""
        for shape in slide.shapes:
            if hasattr(shape, 'text') and ('content' in shape.name.lower() or 'title' in shape.name.lower()):
                logger.debug(f"Found generic content in shape: {shape.name}")
                content += shape.text + " "
        return content.strip()

    def extractKeywords(self, text):
        logger.info("Extracting keywords")
        nlp_method = self.nlpMethod.currentText()
        
        prompt = f"""
        Analyze the following text and extract 4 key concepts that would be most suitable for generating an image. 
        The concepts should capture the essence of the text, even if not directly mentioned. 
        If there is no prompt, please reply in English.
        Provide the 4 keywords or short phrases thinking in the following direction:

        1. Scene: [Overall scene or setting]
        2. Style: [Visual style or artistic approach]
        3. Subject: [Main subject or focus]
        4. Mood: [Emotional tone or atmosphere]
        
        Provide only the 4 keywords or short phrases, separated by commas, without any additional explanation:

        Text: {text}
        """

        if nlp_method == 'Ollama':
            return self.extractKeywordsOllama(prompt)
        elif nlp_method == 'ChatGPT':
            return self.extractKeywordsChatGPT(prompt)
        
    def extractKeywordsOllama(self, prompt):
        logger.info("Extracting keywords using Ollama API")
        ollama_config = self.config.get('nlp', {}).get('ollama', {})
        ollama_url = ollama_config.get('url', "http://localhost:11434/api/generate")
        
        data = {
            "model": ollama_config.get('model', "llama3.1"),
            "prompt": prompt,
            "stream": False
        }
        
        try:
            response = requests.post(ollama_url, json=data)
            response.raise_for_status()
            result = response.json()
            keywords = result['response'].strip()
            logger.info(f"Extracted keywords: {keywords}")
            return keywords
        except requests.RequestException as e:
            logger.error(f"Error in Ollama API call: {e}")
            return ""

    def extractKeywordsChatGPT(self, prompt):
        logger.info("Extracting keywords using ChatGPT API")
        chatgpt_config = self.config.get('nlp', {}).get('chatgpt', {})
        openai.api_key = chatgpt_config.get('api_key', '')
        model = chatgpt_config.get('model', 'gpt-3.5-turbo')  # 从配置文件中获取模型，默认为 'gpt-3.5-turbo'

        try:
            response = openai.ChatCompletion.create(
                model=model,  # 使用从配置文件中读取的模型
                messages=[
                    {"role": "system", "content": "You are a helpful assistant that extracts keywords from text."},
                    {"role": "user", "content": f"Please extract key words from the following text: {prompt}"}
                ]
            )
            keywords = response.choices[0].message['content'].strip()
            logger.info(f"Extracted keywords: {keywords}")
            return keywords
        except openai.error.OpenAIError as e:
            logger.error(f"Error in ChatGPT API call: {e}")
            return ""

    def generateImage(self, keywords, method):
        logger.info(f"Generating image using method: {method}")
        if method == 'Comfyui生成':
            return self.generateImageComfyui(keywords)
        elif method == 'DALL-E':
            return self.generateImageDALLE(keywords)
        elif method == 'Pixabay':
            return self.searchImagePixabay(keywords)
        elif method == 'Unsplash':
            return self.searchImageUnsplash(keywords)

    def generateImageComfyui(self, keywords):
        logger.info(f"Generating image with Comfyui using keywords: {keywords}")
        
        comfyui_config = self.config.get('image_generation', {}).get('comfyui', {})
        comfyui_base_url = comfyui_config.get('base_url', "http://127.0.0.1:8188")
        prompt_url = f"{comfyui_base_url}/prompt"
        
        try:
            with open(comfyui_config.get('workflow_path', './Input/workflow.json'), 'r') as file:
                workflow = json.load(file)
        except FileNotFoundError:
            logger.error("Comfyui workflow file not found")
            return None
        except json.JSONDecodeError:
            logger.error("Invalid JSON in Comfyui workflow file")
            return None

        # 更新工作流中的关键词
        for node in workflow.values():
            if node['class_type'] == 'CLIPTextEncode':
                if node['inputs'].get('text') == '[KEYWORDS]':
                    node['inputs']['text'] = keywords
                elif '[KEYWORDS]' in node['inputs'].get('text', ''):
                    node['inputs']['text'] = node['inputs']['text'].replace('[KEYWORDS]', keywords)

        data = {
            "prompt": workflow,
            "client_id": "ppt_image_generator"
        }
        
        try:
            # 提交工作流
            response = requests.post(prompt_url, json=data)
            response.raise_for_status()
            
            prompt_id = response.json()['prompt_id']
            logger.info(f"Comfyui workflow submitted with ID: {prompt_id}")

            # 轮询等待图像生成完成
            max_attempts = 120  # 最多等待120次，每次5秒，总共10分钟
            for attempt in range(max_attempts):
                time.sleep(5)  # 每5秒检查一次
                history_url = f"{comfyui_base_url}/history/{prompt_id}"
                history_response = requests.get(history_url)
                history_response.raise_for_status()
                history_data = history_response.json()

                logger.debug(f"History data: {json.dumps(history_data, indent=2)}")

                if prompt_id in history_data:
                    if 'outputs' in history_data[prompt_id]:
                        output_images = history_data[prompt_id]['outputs']
                        if output_images:
                            # 假设我们只关心第一个输出图像
                            first_output_key = list(output_images.keys())[0]
                            image_data = output_images[first_output_key]
                            logger.debug(f"Image data: {json.dumps(image_data, indent=2)}")

                            if 'images' in image_data and len(image_data['images']) > 0:
                                first_image = image_data['images'][0]
                                if 'filename' in first_image and 'type' in first_image:
                                    image_url = f"{comfyui_base_url}/view?filename={first_image['filename']}&type={first_image['type']}"
                                    if 'subfolder' in first_image and first_image['subfolder']:
                                        image_url += f"&subfolder={first_image['subfolder']}"
                                    
                                    # 下载并返回图像
                                    img_response = requests.get(image_url)
                                    img_response.raise_for_status()
                                    logger.info("Image successfully retrieved from Comfyui")
                                    return Image.open(BytesIO(img_response.content))
                                else:
                                    logger.error(f"Missing 'filename' or 'type' in image data: {first_image}")
                                    return None
                            else:
                                logger.error(f"No 'images' list in image data: {image_data}")
                                return None
                        else:
                            logger.warning("No image data in Comfyui response")
                            return None
                    elif 'error' in history_data[prompt_id]:
                        logger.error(f"Error in Comfyui image generation: {history_data[prompt_id]['error']}")
                        return None
                
                logger.info(f"Image generation in progress, attempt {attempt + 1}/{max_attempts}")

            logger.warning("Timed out waiting for Comfyui image generation")
            return None

        except requests.RequestException as e:
            logger.error(f"Error in Comfyui API call: {e}")
            return None
        except Exception as e:
            logger.error(f"Unexpected error in generateImageComfyui: {e}")
            return None

    def generateImageDALLE(self, keywords):
        logger.info(f"Generating image with DALL-E using keywords: {keywords}")
        dalle_config = self.config.get('image_generation', {}).get('dalle', {})
        openai.api_key = dalle_config.get('api_key', '')

        try:
            response = openai.Image.create(
                prompt=keywords,
                n=1,
                size="1024x1024"
            )
            image_url = response['data'][0]['url']
            image_response = requests.get(image_url)
            image_response.raise_for_status()
            image = Image.open(BytesIO(image_response.content))
            logger.info("Image generated successfully with DALL-E")
            return image
        except (openai.error.OpenAIError, requests.RequestException) as e:
            logger.error(f"Error generating image with DALL-E: {e}")
            return None

    def searchImagePixabay(self, keywords):
        logger.info(f"Original keywords: {keywords}")
        
        # 提取所有类型的关键词，每类取前三个词
        extracted_keywords = []
        for line in keywords.split('\n'):
            parts = line.split(':', 1)
            if len(parts) == 2:
                keyword = parts[1].strip().rstrip(',')
                keyword_parts = keyword.split()[:2]
                extracted_keywords.extend(keyword_parts)
        
        # 将提取的关键词组合成一个搜索字符串
        search_keywords = ' '.join(extracted_keywords)
        logger.info(f"Extracted keywords for Pixabay search: {search_keywords}")

        pixabay_config = self.config.get('pixabay', {})
        api_key = pixabay_config.get('api_key', '')

        logger.debug(f"Pixabay API key: {api_key[:5]}...{api_key[-5:] if len(api_key) > 10 else ''}")

        url = "https://pixabay.com/api/"

        params = {
            "key": api_key,
            "q": search_keywords,
            "image_type": "photo",
            "per_page": 20,  # 增加返回的图片数量
            "order": "random"  # 使用随机排序
        }

        try:
            response = requests.get(url, params=params)
            response.raise_for_status()
            data = response.json()
            if data["hits"]:
                # 随机选择一张图片
                image_url = random.choice(data["hits"])["largeImageURL"]
                image_response = requests.get(image_url)
                image_response.raise_for_status()
                image = Image.open(BytesIO(image_response.content))
                logger.info(f"Image found successfully on Pixabay for keywords: {search_keywords}")
                return image
            else:
                logger.warning(f"No images found on Pixabay for keywords: {search_keywords}")
                return None
        except requests.RequestException as e:
            logger.error(f"Error searching image on Pixabay: {e}")
            return None
    
    def searchImageUnsplash(self, keywords):
        logger.info(f"Searching image on Unsplash using keywords: {keywords}")
        unsplash_config = self.config.get('image_generation', {}).get('unsplash', {})
        access_key = unsplash_config.get('access_key', '')
        url = "https://api.unsplash.com/search/photos"

        headers = {
            "Authorization": f"Client-ID {access_key}"
        }
        params = {
            "query": keywords,
            "per_page": 1
        }

        try:
            response = requests.get(url, headers=headers, params=params)
            response.raise_for_status()
            data = response.json()
            if data["results"]:
                image_url = data["results"][0]["urls"]["regular"]
                image_response = requests.get(image_url)
                image_response.raise_for_status()
                image = Image.open(BytesIO(image_response.content))
                logger.info("Image found successfully on Unsplash")
                return image
            else:
                logger.warning("No images found on Unsplash")
                return None
        except requests.RequestException as e:
            logger.error(f"Error searching image on Unsplash: {e}")
            return None

    def insertImage(self, slide, shape, image):
        logger.info(f"Inserting new image for shape: {shape.name}")
        # 调整图片大小和裁剪
        placeholder_ratio = shape.width / shape.height
        image_ratio = image.width / image.height
        
        if image_ratio > placeholder_ratio:
            new_width = int(image.height * placeholder_ratio)
            left_margin = (image.width - new_width) / 2
            image = image.crop((left_margin, 0, left_margin + new_width, image.height))
        else:
            new_height = int(image.width / placeholder_ratio)
            top_margin = (image.height - new_height) / 2
            image = image.crop((0, top_margin, image.width, top_margin + new_height))
        
        # 调整图片大小
        image.thumbnail((shape.width, shape.height))
        
        # 保存临时文件
        temp_image_path = 'temp_image.png'
        image.save(temp_image_path)
        logger.debug(f"Temporary image saved as: {temp_image_path}")
        
        # 插入图片
        left = shape.left
        top = shape.top
        width = shape.width
        height = shape.height
        new_picture = slide.shapes.add_picture(temp_image_path, left, top, width, height)
        new_picture.name = shape.name  # 保持原有图片占位符的名字
        logger.info(f"New image inserted into slide with name: {new_picture.name}")
        
        # 将新插入的图片移动到最底层
        slide.shapes._spTree.remove(new_picture._element)
        slide.shapes._spTree.insert(2, new_picture._element)
        logger.info("New image moved to the bottom layer")
        
        # 删除原来的形状
        sp = shape._element
        sp.getparent().remove(sp)
        logger.info("Original shape removed")
        
        # 删除临时文件
        os.remove(temp_image_path)
        logger.debug("Temporary image file removed")

    def all_changes_completed(self):
        # 在这里实现检查所有更改是否完成的逻辑
        # 这个方法应该返回True如果所有更改都完成了，否则返回False
        return True

def main():
    app = QApplication(sys.argv)
    
    # Set the application style to Fusion
    app.setStyle(QStyleFactory.create('Fusion'))
    
    # Optional: Set a dark palette
    dark_palette = QPalette()
    # ... (palette color settings as shown above)
    app.setPalette(dark_palette)
    
    ex = PPTImageGenerator()
    ex.show()
    sys.exit(app.exec_())

if __name__ == '__main__':
    main()