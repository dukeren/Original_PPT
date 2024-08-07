import os
import re
import sys
import random
from PyQt5.QtWidgets import (QApplication, QMainWindow, QWidget, QVBoxLayout, QHBoxLayout, 
                             QLabel, QLineEdit, QPushButton, QFileDialog, QMessageBox)
from PyQt5.QtGui import QIcon, QFont
from PyQt5.QtCore import Qt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import subprocess

def parse_markdown(file_path):
    # 打开并读取Markdown文件
    with open(file_path, 'r', encoding='utf-8') as file:
        content = file.read()
    
    slides = []  # 存储所有幻灯片
    current_slide = {'type': None, 'title': None, 'content': []}  # 当前处理的幻灯片
    
    lines = content.split('\n')  # 将内容分割成行
    subtitle_count = 0  # 四级标题计数
    content_count = 0  # 内容计数
    subcontent_count = 0  # 子内容计数
    last_content_type = None  # 上一个内容的类型
    has_cover = False  # 是否有封面
    chapters = []  # 存储章节标题
    
    def process_content(line, current_content):
        # 检测列表项
        if line.lstrip().startswith(('- ', '* ', '+ ', '1. ', '2. ', '3. ')):
            # 如果是标准列表项,只缩进四个空格
            indent = '    '
            line = indent + line.lstrip()
        elif line.startswith('\t'):  # 检测 TAB 缩进
            # 如果是 TAB 缩进的列表项,替换 TAB 为两个空格
            line = '    ' + line.lstrip('\t')

        # 如果当前内容不为空,添加新行;否则直接返回当前行
        if current_content:
            return current_content + '\n' + line
        else:
            return line

    for line in lines:
        if line.startswith('# '):  # 一级标题:封面
            if current_slide['type']:
                slides.append(current_slide)
            current_slide = {'type': 'cover', 'title': line[2:], 'content': []}
            has_cover = True
            subtitle_count = content_count = subcontent_count = 0
            last_content_type = None
        elif line.startswith('## '):  # 二级标题:章节
            if current_slide['type']:
                slides.append(current_slide)
            current_slide = {'type': 'chapter', 'title': line[3:], 'content': []}
            chapters.append(line[3:])
            subtitle_count = content_count = subcontent_count = 0
            last_content_type = None
        elif line.startswith('### '):  # 三级标题:主要内容
            if current_slide['type']:
                slides.append(current_slide)
            current_slide = {'type': 'substance', 'title': line[4:], 'content': []}
            subtitle_count = content_count = subcontent_count = 0
            last_content_type = None
        elif line.startswith('#### '):  # 四级标题:子标题
            subtitle_count += 1
            current_slide['content'].append((f'subtitle{subtitle_count:02d}', line[5:]))
            subcontent_count = 0
            last_content_type = 'subtitle'
        elif line.strip() == '---':  # 分隔线:翻译幻灯片
            if current_slide['type']:
                slides.append(current_slide)
            current_slide = {'type': 'translate', 'title': None, 'content': []}
            subtitle_count = content_count = subcontent_count = 0
            last_content_type = None
        elif line.strip():  # 非空行:内容
            if current_slide['type'] in ['substance', 'chapter']:
                if last_content_type == 'subtitle' or subcontent_count > 0:
                    subcontent_count += 1
                    content_key = f'subcontent{subtitle_count:02d}'
                    processed_line = process_content(line, '')
                    if subcontent_count == 1:
                        current_slide['content'].append((content_key, processed_line))
                    else:
                        current_content = current_slide['content'][-1]
                        current_slide['content'][-1] = (current_content[0], process_content(processed_line, current_content[1]))
                else:
                    content_count += 1
                    content_key = f'content{content_count:02d}'
                    if content_count == 1 or last_content_type != 'content':
                        current_slide['content'].append((content_key, process_content(line, '')))
                    else:
                        current_content = current_slide['content'][-1]
                        current_slide['content'][-1] = (current_content[0], process_content(line, current_content[1]))
            else:
                content_count += 1
                current_slide['content'].append((f'content{content_count:02d}', process_content(line, '')))
            last_content_type = 'content'
    
    # 添加最后一个幻灯片
    if current_slide['type']:
        slides.append(current_slide)
    
    # 如果有封面,在封面后添加目录
    if has_cover:
        toc_slide = {'type': 'toc', 'title': '目录', 'content': [('content01', '\n'.join(chapters))]}
        slides.insert(1, toc_slide)
    
    # 打印检测内容
    print("解析结果:")
    for i, slide in enumerate(slides):
        print(f"幻灯片 {i+1}:")
        print(f"  类型: {slide['type']}")
        print(f"  标题: {slide['title']}")
        print("  内容:")
        for content_type, text in slide['content']:
            print(f"    - [{content_type}] {text}")
        print()
    
    return slides

def rename_placeholders(slide, master_slide):
    for shape in slide.placeholders:
        if hasattr(shape, 'placeholder_format'):
            master_placeholder = master_slide.placeholders.get(shape.placeholder_format.idx)
            if master_placeholder:
                shape.name = master_placeholder.name

def create_pptx(slides, template_file, output_file):
    prs = Presentation(template_file)
    
    print("Available layouts in the template:")
    for layout in prs.slide_layouts:
        print(f"  - {layout.name}")
    
    for slide_index, slide in enumerate(slides, 1):
        slide_type = slide['type']
        print(f"\nProcessing slide {slide_index}, type: {slide_type}")
        
        if slide_type.lower() == 'substance':
            subtitle_count = sum(1 for content_type, _ in slide['content'] if content_type.lower().startswith('subtitle'))
            pattern = f"substance_{subtitle_count:02d}"
            print(f"Substance slide detected. Subtitle count: {subtitle_count}")
            print(f"Searching for layout with pattern: {pattern}")
            matching_layouts = [layout for layout in prs.slide_layouts if pattern.lower() in layout.name.lower()]
        else:
            print(f"Non-substance slide. Searching for layout with type: {slide_type}")
            matching_layouts = [layout for layout in prs.slide_layouts if slide_type.lower() in layout.name.lower()]
        
        print(f"Found {len(matching_layouts)} matching layouts:")
        for layout in matching_layouts:
            print(f"  - {layout.name}")
        
        if matching_layouts:
            slide_layout = random.choice(matching_layouts)
            print(f"Selected layout: {slide_layout.name}")
        else:
            print(f"Warning: No layout found for slide type '{slide_type}'. Using default layout.")
            slide_layout = prs.slide_layouts[0]  # 使用默认布局
            print(f"Selected default layout: {slide_layout.name}")

        new_slide = prs.slides.add_slide(slide_layout)
        rename_placeholders(new_slide, slide_layout)  # 重命名占位符
        
        print(f"\n--- Slide {slide_index} ({slide['type']}) ---")
        print("Available placeholders after renaming:")
        for shape in new_slide.placeholders:
            print(f"  - {shape.name} (index: {shape.placeholder_format.idx})")
        
        # 设置标题
        if slide['title']:
            title_placeholder = find_placeholder(new_slide, 'title')
            if title_placeholder:
                original_font = title_placeholder.text_frame.paragraphs[0].font
                original_size = original_font.size
                original_name = original_font.name
                original_color = original_font.color.rgb if hasattr(original_font.color, 'rgb') else None
                original_bold = original_font.bold
                original_italic = original_font.italic

                title_placeholder.text = slide['title']
                
                # 重新应用原有的字体设置
                new_font = title_placeholder.text_frame.paragraphs[0].font
                new_font.size = original_size
                new_font.name = original_name
                if original_color:
                    new_font.color.rgb = original_color
                new_font.bold = original_bold
                new_font.italic = original_italic

                print(f"Title set: {slide['title']}")
            else:
                print("WARNING: No title placeholder found")
        
        # 处理内容
        for content_type, text in slide['content']:
            placeholder = find_placeholder(new_slide, content_type)
            if placeholder:
                tf = placeholder.text_frame
                if tf.paragraphs:
                    original_paragraph = tf.paragraphs[0]
                    original_font = original_paragraph.font
                    original_size = original_font.size
                    original_name = original_font.name
                    original_color = original_font.color.rgb if hasattr(original_font.color, 'rgb') else None
                    original_bold = original_font.bold
                    original_italic = original_font.italic
                    original_alignment = original_paragraph.alignment

                    p = original_paragraph
                    p.text = text
                else:
                    p = tf.add_paragraph()
                    p.text = text
                    original_size = Pt(18)  # 默认大小
                    original_name = 'Calibri'  # 默认字体
                    original_color = None
                    original_bold = None
                    original_italic = None
                    original_alignment = PP_ALIGN.LEFT

                # 重新应用原有的字体和段落设置
                new_font = p.font
                new_font.size = original_size
                new_font.name = original_name
                if original_color:
                    new_font.color.rgb = original_color
                if original_bold is not None:
                    new_font.bold = original_bold
                if original_italic is not None:
                    new_font.italic = original_italic
                p.alignment = original_alignment

                print(f"Content replaced in {content_type}: {text[:30]}...")
            else:
                print(f"WARNING: No placeholder found for {content_type}")
    
    prs.save(output_file)
    print(f"\nPresentation saved as {output_file}")

def find_placeholder(slide, content_type):
    # 对于标题，直接查找 'Title' 占位符
    if content_type.lower() == 'title':
        for shape in slide.placeholders:
            if shape.name.lower() == 'title':
                print(f"Found title placeholder: {shape.name}")
                return shape
        return None

    # 提取内容类型和序号
    match = re.match(r'(\w+)(\d+)', content_type)
    if match:
        base_type, number = match.groups()
        # 查找完全匹配的占位符
        for shape in slide.placeholders:
            if shape.name.lower() == content_type.lower():
                print(f"Found exact match for {content_type}: {shape.name}")
                return shape
        # 如果没有找到完全匹配的，查找基本类型匹配的
        for shape in slide.placeholders:
            if shape.name.lower().startswith(base_type.lower()):
                print(f"Found base match for {content_type}: {shape.name}")
                return shape
    print(f"No placeholder found for {content_type}")
    return None

class PPTGeneratorGUI(QMainWindow):
    def __init__(self):
        super().__init__()
        self.initUI()

    def initUI(self):
        self.setWindowTitle("PPT生成 By 渡客")
        self.setGeometry(100, 100, 700, 300)  # 稍微增加窗口大小

        # Set icon
        icon_path = os.path.join(os.path.dirname(__file__), "Image", "logo.png")
        if os.path.exists(icon_path):
            self.setWindowIcon(QIcon(icon_path))

        # Main widget and layout
        central_widget = QWidget(self)
        self.setCentralWidget(central_widget)
        main_layout = QVBoxLayout(central_widget)
        main_layout.setContentsMargins(20, 50, 20, 20)  # 设置主布局的边距
        main_layout.setSpacing(15)  # 设置主布局中小部件之间的间距

        # Markdown file
        markdown_layout = QHBoxLayout()
        markdown_layout.setSpacing(10)  # 设置水平布局中小部件之间的间距
        markdown_label = QLabel("Markdown文件:")
        self.markdown_entry = QLineEdit()
        markdown_button = QPushButton("浏览")
        markdown_button.clicked.connect(self.browse_markdown)
        markdown_layout.addWidget(markdown_label)
        markdown_layout.addWidget(self.markdown_entry)
        markdown_layout.addWidget(markdown_button)
        main_layout.addLayout(markdown_layout)

        # Template file
        template_layout = QHBoxLayout()
        template_layout.setSpacing(10)
        template_label = QLabel("PPT 参考模板:")
        self.template_entry = QLineEdit()
        template_button = QPushButton("浏览")
        template_button.clicked.connect(self.browse_template)
        template_layout.addWidget(template_label)
        template_layout.addWidget(self.template_entry)
        template_layout.addWidget(template_button)
        main_layout.addLayout(template_layout)

        # Output directory
        output_layout = QHBoxLayout()
        output_layout.setSpacing(10)
        output_label = QLabel("文件保存路径:")
        self.output_entry = QLineEdit()
        output_button = QPushButton("浏览")
        output_button.clicked.connect(self.browse_output)
        output_layout.addWidget(output_label)
        output_layout.addWidget(self.output_entry)
        output_layout.addWidget(output_button)
        main_layout.addLayout(output_layout)

        # Add some vertical space
        main_layout.addSpacing(20)

        # Generate button
        generate_button = QPushButton("一键生成 PPT")
        generate_button.clicked.connect(self.generate_ppt)
        generate_button.setFixedSize(200, 40)  # 设置按钮大小
        main_layout.addWidget(generate_button, alignment=Qt.AlignCenter)

        # Add some vertical space
        main_layout.addSpacing(10)

        # Help link
        help_link = QLabel("使用说明")
        help_link.setStyleSheet("color: blue; text-decoration: underline;")
        help_link.setCursor(Qt.PointingHandCursor)
        help_link.mousePressEvent = self.open_readme
        main_layout.addWidget(help_link, alignment=Qt.AlignRight)

        # Set font
        font = QFont("Helvetica", 12)
        self.setFont(font)

        # Set style sheet for better looking buttons and entries
        self.setStyleSheet("""
            QPushButton {
                padding: 5px 10px;
                background-color: #f0f0f0;
                border: 1px solid #ccc;
                border-radius: 3px;
            }
            QPushButton:hover {
                background-color: #e0e0e0;
            }
            QLineEdit {
                padding: 5px;
                border: 1px solid #ccc;
                border-radius: 3px;
            }
        """)

    def browse_markdown(self):
        filename, _ = QFileDialog.getOpenFileName(self, "选择Markdown文件", "", "Markdown Files (*.md)")
        if filename:
            self.markdown_entry.setText(filename)

    def browse_template(self):
        filename, _ = QFileDialog.getOpenFileName(self, "选择PPT模板", "", "PowerPoint Files (*.pptx)")
        if filename:
            self.template_entry.setText(filename)

    def browse_output(self):
        directory = QFileDialog.getExistingDirectory(self, "选择保存路径")
        if directory:
            self.output_entry.setText(directory)

    def generate_ppt(self):
        markdown_file = self.markdown_entry.text()
        template_file = self.template_entry.text()
        output_directory = self.output_entry.text()

        # Check and process inputs
        if not markdown_file:
            markdown_file = 'Input/input.md'
            if not os.path.exists(markdown_file):
                QMessageBox.critical(self, "错误", f"默认的Markdown文件 '{markdown_file}' 未找到。")
                return
            self.markdown_entry.setText(markdown_file)

        if not template_file:
            template_file = 'Model_PPT/Model.pptx'
            if not os.path.exists(template_file):
                QMessageBox.critical(self, "错误", f"默认的模板文件 '{template_file}' 未找到。")
                return
            self.template_entry.setText(template_file)

        if not output_directory:
            output_directory = 'Outfile'
            if not os.path.exists(output_directory):
                os.makedirs(output_directory)
            self.output_entry.setText(output_directory)

        try:
            # Generate a default output filename
            output_filename = "output.pptx"
            output_path = os.path.join(output_directory, output_filename)

            # If file already exists, add a number to the filename
            counter = 1
            while os.path.exists(output_path):
                output_filename = f"output_{counter}.pptx"
                output_path = os.path.join(output_directory, output_filename)
                counter += 1

            slides = parse_markdown(markdown_file)
            create_pptx(slides, template_file, output_path)
            QMessageBox.information(self, "成功", f"PPT生成成功: {output_path}")
            
            # 自动打开生成文件的目录
            self.open_output_directory(output_directory)
        except Exception as e:
            QMessageBox.critical(self, "错误", f"发生错误: {str(e)}")

    def open_output_directory(self, path):
        if os.path.exists(path):
            if os.name == 'nt':  # Windows
                os.startfile(path)
            elif os.name == 'posix':  # macOS and Linux
                try:
                    subprocess.call(['open', path])  # macOS
                except:
                    subprocess.call(['xdg-open', path])  # Linux
        else:
            QMessageBox.critical(self, "错误", f"目录未找到: {path}")

    def open_readme(self, event):
        if getattr(sys, 'frozen', False):
            # If it's a packaged executable
            application_path = sys._MEIPASS
        else:
            # If it's a script
            application_path = os.path.dirname(os.path.abspath(__file__))
        
        readme_path = os.path.join(application_path, "readme.md")
        
        if os.path.exists(readme_path):
            if os.name == 'nt':  # Windows
                os.startfile(readme_path)
            elif os.name == 'posix':  # macOS and Linux
                try:
                    subprocess.call(['open', readme_path])  # macOS
                except:
                    subprocess.call(['xdg-open', readme_path])  # Linux
        else:
            QMessageBox.critical(self, "错误", f"readme.md 文件未在 {readme_path} 找到")

if __name__ == "__main__":
    app = QApplication(sys.argv)
    app.setStyle("Fusion")  # Set Fusion style
    gui = PPTGeneratorGUI()
    gui.show()
    sys.exit(app.exec_())